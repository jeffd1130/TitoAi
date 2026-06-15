#!/usr/bin/env python3
"""Generate Tito AI Niche & Growth Plan PPTX — June 2026."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

OUT = "/Users/jeff/Documents/Claude/TItoAi/content/strategy/TitoAI-Niche-Growth-Plan-June2026.pptx"

# Brand colors
NAVY   = RGBColor(0x0A, 0x0F, 0x1E)
GOLD   = RGBColor(0xF5, 0x9E, 0x0B)
TEAL   = RGBColor(0x0D, 0x94, 0x88)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x9C, 0xA3, 0xAF)
LGRAY  = RGBColor(0xE5, 0xE7, 0xEB)
DARK   = RGBColor(0x11, 0x18, 0x27)
RED    = RGBColor(0xDC, 0x26, 0x26)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]  # completely blank


def add_slide():
    return prs.slides.add_slide(blank)


def bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def txbox(slide, text, l, t, w, h,
          size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          wrap=True, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"
    return tb


def label(slide, text, l, t, w, h, color=GOLD):
    """Small all-caps label."""
    txbox(slide, text.upper(), l, t, w, h, size=9, bold=True, color=color)


def gold_bar(slide, t=Inches(0.08)):
    rect(slide, 0, 0, W, Pt(6), GOLD)


def footer_bar(slide):
    rect(slide, 0, H - Inches(0.45), W, Inches(0.45), DARK)
    txbox(slide, "Tito AI @TitoAIPH  ·  Niche & Growth Plan  ·  June 2026",
          Inches(0.4), H - Inches(0.38), Inches(8), Inches(0.32),
          size=8, color=GRAY)
    txbox(slide, "Confidential — Internal use only",
          Inches(9), H - Inches(0.38), Inches(4), Inches(0.32),
          size=8, color=GRAY, align=PP_ALIGN.RIGHT)


def slide_number(slide, n):
    txbox(slide, str(n), W - Inches(0.55), H - Inches(0.38),
          Inches(0.4), Inches(0.32), size=8, color=GOLD, align=PP_ALIGN.RIGHT)


def section_header(slide, text):
    rect(slide, 0, Inches(1.15), Inches(0.06), Inches(0.55), GOLD)
    txbox(slide, text, Inches(0.2), Inches(1.1), Inches(12), Inches(0.65),
          size=28, bold=True, color=WHITE)
    rect(slide, Inches(0.2), Inches(1.82), Inches(12.9), Pt(2), GOLD)


# ── SLIDE 1 — COVER ─────────────────────────────────────────────────────────
s = add_slide(); bg(s); n = 1
gold_bar(s)
rect(s, 0, H - Inches(2.2), W, Inches(2.2), DARK)
txbox(s, "CONFIDENTIAL · INTERNAL STRATEGY", Inches(0.6), Inches(0.5),
      Inches(10), Inches(0.4), size=9, bold=True, color=GOLD)
txbox(s, "Tito AI", Inches(0.6), Inches(1.1), Inches(12), Inches(1.1),
      size=60, bold=True, color=WHITE)
txbox(s, "@TitoAIPH", Inches(0.6), Inches(2.1), Inches(12), Inches(0.7),
      size=36, bold=True, color=GOLD)
txbox(s, "Niche & Growth Plan", Inches(0.6), Inches(2.85), Inches(12), Inches(0.6),
      size=24, bold=False, color=LGRAY)
rect(s, Inches(0.6), Inches(3.55), Inches(0.8), Pt(4), GOLD)
# tagline box
tl = rect(s, Inches(0.6), Inches(3.8), Inches(11.5), Inches(1.1), RGBColor(0x1E, 0x2A, 0x3A))
txbox(s, '"The warm, relatable Tito who teaches everyday Filipinos to use AI for free\n— before it replaces their job."',
      Inches(0.85), Inches(3.85), Inches(11.1), Inches(1.0),
      size=14, italic=True, color=RGBColor(0xFC, 0xD3, 0x4D))
# meta row
for i, (lbl, val) in enumerate([
    ("VERSION", "June 2026 — Post-Audit"),
    ("PREPARED BY", "Jeff de las Armas"),
    ("CHANNEL", "@tito.aiph · TikTok · IG · FB"),
]):
    x = Inches(0.6 + i * 4.3)
    txbox(s, lbl, x, H - Inches(1.95), Inches(4), Inches(0.25), size=8, bold=True, color=GRAY)
    txbox(s, val, x, H - Inches(1.65), Inches(4), Inches(0.35), size=13, color=WHITE)
footer_bar(s); slide_number(s, n)


# ── SLIDE 2 — MARKET OPPORTUNITY ────────────────────────────────────────────
s = add_slide(); bg(s); n = 2
gold_bar(s)
section_header(s, "Market Opportunity")

stats = [
    ("12.7M", "Filipino workers exposed to GenAI\nHighest in ASEAN (ILO 2026)"),
    ("42.4%", "Filipino internet users use ChatGPT monthly\n6th globally (Radar PH 2026)"),
    ("1.9M", "BPO employees who need to upskill now\n(BSP 2025)"),
    ("1.5M+", "Filipino freelancers underserved\nby AI education (PIDS)"),
    ("14.9%", "Only — of small firms currently\nuse AI tools (PH AI Report 2025)"),
    ("83%", "Filipino students use AI\nParents haven't caught up (BW 2026)"),
]
cols, rows = 3, 2
cw, ch = Inches(4.2), Inches(1.5)
for i, (num, desc) in enumerate(stats):
    col, row = i % cols, i // cols
    x = Inches(0.3 + col * 4.38)
    y = Inches(2.05 + row * 1.65)
    rect(s, x, y, cw, ch, DARK)
    rect(s, x, y, Pt(4), ch, GOLD)
    txbox(s, num, x + Inches(0.12), y + Inches(0.12), Inches(3.8), Inches(0.65),
          size=30, bold=True, color=GOLD)
    txbox(s, desc, x + Inches(0.12), y + Inches(0.7), Inches(3.8), Inches(0.72),
          size=10, color=LGRAY)

txbox(s, "The gap: 12.7M Filipinos exposed to AI disruption — almost none have a relatable Taglish-speaking guide. Tito AI is the first mover.",
      Inches(0.3), H - Inches(0.85), Inches(12.5), Inches(0.38), size=10, color=TEAL, italic=True)
footer_bar(s); slide_number(s, n)


# ── SLIDE 3 — NICHE & MOAT ──────────────────────────────────────────────────
s = add_slide(); bg(s); n = 3
gold_bar(s)
section_header(s, "Niche & Competitive Moat")

txbox(s, "Target Audience", Inches(0.3), Inches(2.0), Inches(6), Inches(0.3),
      size=10, bold=True, color=GOLD)
audience = ["Freelancers (1.5M+)", "Guro (teachers)", "BPO workers afraid of being replaced",
            "Nanays / Tatays", "Small business owners"]
for i, item in enumerate(audience):
    rect(s, Inches(0.3), Inches(2.4 + i * 0.52), Inches(0.06), Inches(0.32), TEAL)
    txbox(s, item, Inches(0.5), Inches(2.38 + i * 0.52), Inches(5.5), Inches(0.38),
          size=13, color=WHITE)

txbox(s, "Tito AI's Moat", Inches(7.0), Inches(2.0), Inches(6), Inches(0.3),
      size=10, bold=True, color=GOLD)
moat = [
    ("Taglish", "Barrier to entry for foreign creators"),
    ('"Libre lahat"', "Removes the #1 objection — cost"),
    ("Warm Tito brand", "No one else owns this in PH"),
    ("First mover", 'Algorithm surfaces "AI tutorial Tagalog"'),
]
for i, (title, desc) in enumerate(moat):
    y = Inches(2.4 + i * 0.9)
    rect(s, Inches(7.0), y, Inches(5.9), Inches(0.75), DARK)
    txbox(s, title, Inches(7.15), y + Inches(0.05), Inches(5.5), Inches(0.28),
          size=12, bold=True, color=GOLD)
    txbox(s, desc, Inches(7.15), y + Inches(0.32), Inches(5.5), Inches(0.28),
          size=10, color=LGRAY)

rect(s, Inches(6.7), Inches(1.9), Pt(1.5), Inches(5.5), GRAY)
footer_bar(s); slide_number(s, n)


# ── SLIDE 4 — PLATFORM STRATEGY ─────────────────────────────────────────────
s = add_slide(); bg(s); n = 4
gold_bar(s)
section_header(s, "Platform Strategy")

platforms = [
    ("🥇 PRIMARY", "TikTok", "Discovery engine — viral reach",
     "Algorithm-favored education niche. Zero ad spend needed to go viral. 4×/week minimum.", GOLD),
    ("🥈 SECONDARY", "Instagram", "Authority + depth",
     "Reels + Carousels drive saves. 25–34 urban professionals. Best for follower conversion.", RGBColor(0xE8, 0x79, 0xF9)),
    ("🥉 SUPPORT", "Facebook", "Paid boost + community",
     "Cheapest PH CPM (65–75% cheaper than US). Best for boosting proven content.", RGBColor(0x93, 0xC5, 0xFD)),
]
for i, (badge, name, role, desc, col) in enumerate(platforms):
    x = Inches(0.3 + i * 4.35)
    y = Inches(2.1)
    rect(s, x, y, Inches(4.1), Inches(4.0), DARK)
    rect(s, x, y, Inches(4.1), Pt(4), col)
    txbox(s, badge, x + Inches(0.15), y + Inches(0.18), Inches(3.8), Inches(0.3),
          size=9, bold=True, color=col)
    txbox(s, name, x + Inches(0.15), y + Inches(0.52), Inches(3.8), Inches(0.55),
          size=22, bold=True, color=WHITE)
    txbox(s, role, x + Inches(0.15), y + Inches(1.05), Inches(3.8), Inches(0.35),
          size=11, color=col)
    rect(s, x + Inches(0.15), y + Inches(1.45), Inches(3.6), Pt(1), RGBColor(0x1E, 0x2A, 0x3A))
    txbox(s, desc, x + Inches(0.15), y + Inches(1.6), Inches(3.8), Inches(1.3),
          size=10, color=LGRAY)

rect(s, 0, H - Inches(1.0), W, Inches(0.55), RGBColor(0x7F, 0x1D, 0x1D))
txbox(s, "⚠️  Critical: Convert both Facebook personal profiles into one Facebook Page before running any ads.",
      Inches(0.3), H - Inches(0.98), Inches(12.5), Inches(0.45), size=10, bold=True, color=WHITE)
footer_bar(s); slide_number(s, n)


# ── SLIDE 5 — PROFILE BIOS ──────────────────────────────────────────────────
s = add_slide(); bg(s); n = 5
gold_bar(s)
section_header(s, "Profile Bios — Ready to Paste")

bios = [
    ("TIKTOK · @tito.aiph · 71 chars",
     "Libre AI lessons para sa Pilipinas 🇵🇭\nClaude · Gemini · Walang bayad"),
    ("INSTAGRAM · @tito.aiph · Max 150 chars",
     "AI lessons nang libre para sa lahat 🇵🇭\nClaude · Gemini — walang bayad\nPara sa guro · freelancer · nanay · negosyante\nMga Pamangkin, sama-sama tayong matuto 🤝"),
    ("FACEBOOK PAGE · Short Description",
     "Nagtuturo ng AI nang libre para sa lahat ng Pilipino.\nClaude, Gemini — hakbang-hakbang, Taglish, walang bayad.\nPara sa guro, freelancer, nanay, at lahat na nag-aalala sa AI. 🇵🇭"),
]
for i, (lbl, bio) in enumerate(bios):
    x = Inches(0.3 + i * 4.35)
    y = Inches(2.1)
    rect(s, x, y, Inches(4.1), Inches(4.1), DARK)
    rect(s, x, y, Inches(4.1), Pt(3), TEAL)
    txbox(s, lbl, x + Inches(0.15), y + Inches(0.14), Inches(3.8), Inches(0.32),
          size=8, bold=True, color=TEAL)
    rect(s, x + Inches(0.15), y + Inches(0.5), Inches(3.75), Pt(1), RGBColor(0x1E, 0x2A, 0x3A))
    txbox(s, bio, x + Inches(0.15), y + Inches(0.65), Inches(3.8), Inches(3.2),
          size=11, color=WHITE)

footer_bar(s); slide_number(s, n)


# ── SLIDE 6 — POSTING SCHEDULE ──────────────────────────────────────────────
s = add_slide(); bg(s); n = 6
gold_bar(s)
section_header(s, "Weekly Posting Schedule")

schedule = [
    ("MONDAY",    "TikTok + IG",    "AI Tip Reel · 30–60s",        "Quick tool tip — talking head"),
    ("WEDNESDAY", "TikTok + IG",    "Demo Reel + Carousel · 60–90s","Screen record + Taglish voiceover"),
    ("THURSDAY",  "TikTok",         "Trend reaction / AI news",     "Commentary + screen record"),
    ("SATURDAY",  "Instagram",      "Reel",                         "Before/after productivity story"),
    ("ANYTIME",   "Facebook",       "Cross-post",                   "Wednesday carousel to FB Page"),
]
col_w = [Inches(1.6), Inches(2.2), Inches(3.5), Inches(5.4)]
col_x = [Inches(0.3), Inches(1.95), Inches(4.2), Inches(7.75)]
hdrs  = ["DAY", "PLATFORM", "FORMAT", "CONTENT TYPE"]

for j, (hdr, x) in enumerate(zip(hdrs, col_x)):
    rect(s, x, Inches(2.05), col_w[j], Inches(0.38), DARK)
    txbox(s, hdr, x + Inches(0.08), Inches(2.1), col_w[j], Inches(0.3),
          size=9, bold=True, color=GOLD)

for i, row in enumerate(schedule):
    y = Inches(2.5 + i * 0.68)
    row_bg = RGBColor(0x0D, 0x16, 0x26) if i % 2 == 0 else DARK
    for j, (cell, x) in enumerate(zip(row, col_x)):
        rect(s, x, y, col_w[j], Inches(0.62), row_bg)
        col = GOLD if j == 0 else (TEAL if j == 1 else WHITE)
        txbox(s, cell, x + Inches(0.08), y + Inches(0.1), col_w[j] - Inches(0.1), Inches(0.48),
              size=11, bold=(j == 0), color=col)

txbox(s, "Cross-posting rule: All TikToks → upload natively to Instagram Reels. Never use TikTok cross-post link — native uploads get more reach.",
      Inches(0.3), H - Inches(0.85), Inches(12.5), Inches(0.38), size=9, italic=True, color=GRAY)
footer_bar(s); slide_number(s, n)


# ── SLIDE 7 — CONTENT PILLARS & HOOK ────────────────────────────────────────
s = add_slide(); bg(s); n = 7
gold_bar(s)
section_header(s, "Content Pillars & Hook Formula")

pillars = [
    ("PILLAR 1", "Tool Demo", "Mon / Wed", GOLD,
     "Highest save rate.\n\"Paano gamitin ang Claude para sa negosyo mo — libre\""),
    ("PILLAR 2", "Before / After", "Saturday", TEAL,
     "Emotional, shareable.\n\"3 oras ng trabaho → 5 minuto gamit ang AI\""),
    ("PILLAR 3", "Fear → Reassurance", "Evergreen", RGBColor(0xE8, 0x79, 0xF9),
     "High anxiety = high engagement.\n\"Hindi ka papalitan ng AI. Papalitan ka ng taong gumagamit.\""),
    ("PILLAR 4", "Everyday Filipino Life", "Friday", RGBColor(0x34, 0xD3, 0x99),
     "Warmth + identity.\nGuro, nanay, negosyante stories."),
]
for i, (tag, name, day, col, desc) in enumerate(pillars):
    x = Inches(0.3 + i * 3.25)
    y = Inches(2.1)
    rect(s, x, y, Inches(3.0), Inches(2.8), DARK)
    rect(s, x, y, Inches(3.0), Pt(3), col)
    txbox(s, tag, x + Inches(0.12), y + Inches(0.12), Inches(2.7), Inches(0.25),
          size=8, bold=True, color=col)
    txbox(s, name, x + Inches(0.12), y + Inches(0.42), Inches(2.7), Inches(0.45),
          size=14, bold=True, color=WHITE)
    txbox(s, day, x + Inches(0.12), y + Inches(0.88), Inches(2.7), Inches(0.28),
          size=9, color=col)
    txbox(s, desc, x + Inches(0.12), y + Inches(1.25), Inches(2.7), Inches(1.4),
          size=9, color=LGRAY)

txbox(s, "Hook Formula — Every video opens in second 1 with ONE of:", Inches(0.3), Inches(5.15),
      Inches(12.5), Inches(0.3), size=10, bold=True, color=GOLD)
hooks = [
    ("Fear", "\"Kung BPO worker ka, pakinggan mo 'to.\""),
    ("Result", "\"30 seconds lang. Tingnan mo 'to.\" → show immediately"),
    ("Identity", "\"May negosyo ka at wala kang marketing budget?\""),
    ("Curiosity", "\"Ito ang tanong na hindi mo pa naitatanong sa AI.\""),
]
for i, (t, h) in enumerate(hooks):
    x = Inches(0.3 + i * 3.25)
    rect(s, x, Inches(5.5), Inches(3.0), Inches(1.4), RGBColor(0x0D, 0x16, 0x26))
    txbox(s, t.upper(), x + Inches(0.12), Inches(5.58), Inches(2.7), Inches(0.25),
          size=8, bold=True, color=TEAL)
    txbox(s, h, x + Inches(0.12), Inches(5.88), Inches(2.7), Inches(0.85),
          size=9, italic=True, color=LGRAY)
footer_bar(s); slide_number(s, n)


# ── SLIDE 8 — BOOST STRATEGY ────────────────────────────────────────────────
s = add_slide(); bg(s); n = 8
gold_bar(s)
section_header(s, "Boost Strategy")

txbox(s, "Rule: Only boost posts with strong organic performance — proven signal = lower CPM.",
      Inches(0.3), Inches(2.1), Inches(12.5), Inches(0.38), size=12, color=WHITE)

# process
steps = ["Post organically", "Wait 24–48 hours", "50+ saves OR 200+ shares?", "Boost it ✓"]
for i, step in enumerate(steps):
    x = Inches(0.3 + i * 3.26)
    circ_size = Inches(0.5)
    cx = x + Inches(1.3)
    cy = Inches(2.72)
    rect(s, cx, cy, circ_size, circ_size, GOLD)
    txbox(s, str(i + 1), cx + Inches(0.13), cy + Inches(0.04), Inches(0.28), Inches(0.35),
          size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    txbox(s, step, x, cy + Inches(0.6), Inches(3.0), Inches(0.4),
          size=11, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 3:
        txbox(s, "→", x + Inches(2.85), cy + Inches(0.08), Inches(0.5), Inches(0.38),
              size=18, color=GOLD, align=PP_ALIGN.CENTER)

# budget tiers
txbox(s, "Budget Tiers", Inches(0.3), Inches(3.75), Inches(6), Inches(0.3),
      size=10, bold=True, color=GOLD)
tiers = [
    ("TEST", "Php 500–1,000 / day", "First boost — prove the format"),
    ("SCALE", "Php 2,000–3,000 / day", "Format proven, engagement positive"),
]
for i, (stage, budget, when) in enumerate(tiers):
    y = Inches(4.1 + i * 0.82)
    rect(s, Inches(0.3), y, Inches(6.2), Inches(0.7), DARK)
    txbox(s, stage, Inches(0.45), y + Inches(0.08), Inches(1.2), Inches(0.28),
          size=10, bold=True, color=GOLD)
    txbox(s, budget, Inches(1.75), y + Inches(0.08), Inches(2.5), Inches(0.28),
          size=12, bold=True, color=WHITE)
    txbox(s, when, Inches(1.75), y + Inches(0.38), Inches(4.5), Inches(0.25),
          size=9, color=GRAY)

# targeting
txbox(s, "Targeting", Inches(7.0), Inches(3.75), Inches(6), Inches(0.3),
      size=10, bold=True, color=GOLD)
targeting = [
    "📍 Location: Philippines",
    "👤 Age: 25–40",
    "💡 Interests: technology, online business, productivity, freelancing",
    "🚀 Best boost candidate: 60–90s screen record + Taglish voiceover",
    "📱 First boost platform: Facebook Page (convert profiles first)",
]
for i, item in enumerate(targeting):
    txbox(s, item, Inches(7.0), Inches(4.18 + i * 0.48), Inches(6.0), Inches(0.4),
          size=10, color=LGRAY)
footer_bar(s); slide_number(s, n)


# ── SLIDE 9 — KPIs ──────────────────────────────────────────────────────────
s = add_slide(); bg(s); n = 9
gold_bar(s)
section_header(s, "KPIs to Track Weekly")

kpis = [
    ("Followers — TikTok",   "100–500",  "2,000+"),
    ("Avg video views",       "500–1,000","5,000+"),
    ("Watch completion rate", ">40%",     ">55%"),
    ("Saves per post",        "20+",      "100+"),
    ("Comments per video",    "10+",      "30+"),
    ("Profile visits",        "100+",     "500+"),
]
hdrs2 = ["METRIC", "MONTH 1 TARGET", "MONTH 3 TARGET"]
hx    = [Inches(0.3), Inches(6.5), Inches(10.0)]
hw    = [Inches(6.0), Inches(3.3), Inches(3.1)]

for j, (h, x, w) in enumerate(zip(hdrs2, hx, hw)):
    rect(s, x, Inches(2.1), w, Inches(0.42), RGBColor(0x06, 0x5F, 0x46))
    txbox(s, h, x + Inches(0.1), Inches(2.17), w - Inches(0.1), Inches(0.3),
          size=9, bold=True, color=WHITE)

for i, (metric, m1, m3) in enumerate(kpis):
    y = Inches(2.6 + i * 0.62)
    row_bg = RGBColor(0x0D, 0x16, 0x26) if i % 2 == 0 else DARK
    for j, (val, x, w) in enumerate(zip([metric, m1, m3], hx, hw)):
        rect(s, x, y, w, Inches(0.55), row_bg)
        col = WHITE if j == 0 else TEAL if j == 1 else GOLD
        txbox(s, val, x + Inches(0.1), y + Inches(0.1), w - Inches(0.1), Inches(0.38),
              size=12, bold=(j > 0), color=col)

txbox(s, "Pull weekly: TikTok Creator Center → Analytics  ·  Instagram Professional Dashboard → Insights  ·  Share screenshots with Jeff",
      Inches(0.3), H - Inches(0.85), Inches(12.5), Inches(0.38), size=9, italic=True, color=GRAY)
footer_bar(s); slide_number(s, n)


# ── SLIDE 10 — ACTION CHECKLIST ─────────────────────────────────────────────
s = add_slide(); bg(s); n = 10
gold_bar(s)
section_header(s, "Immediate Action Checklist")

col1 = [
    ("PROFILE FIXES — THIS WEEK", GOLD, [
        "Convert both FB personal profiles into one Facebook Page",
        "Paste new bios on all 4 platforms",
        "Same profile photo across all platforms",
        "Add link-in-bio (IG + TikTok) → GitHub hub",
        "Pin best-performing post to top of each profile",
    ]),
]
col2 = [
    ("CONTENT SETUP — THIS MONTH", TEAL, [
        "Film \"5 Libreng AI Tools\" Reel — priority boost candidate",
        "Create carousel: 5 AI tools for Filipino freelancers",
        "Film 2 more tool demo Reels (1 per week minimum)",
        "Cross-post all TikToks to Instagram Reels natively",
    ]),
    ("AUTOMATION — WEEKLY CADENCE", RGBColor(0xE8, 0x79, 0xF9), [
        "Monday: run weekly content brief in Claude Code",
        "Friday: run trend scan for following week planning",
        "Monthly: competitor scan to update benchmarks",
    ]),
]

for i, (title, col, items) in enumerate(col1):
    x = Inches(0.3)
    y = Inches(2.1)
    txbox(s, title, x, y, Inches(6.2), Inches(0.3), size=9, bold=True, color=col)
    for j, item in enumerate(items):
        iy = y + Inches(0.45 + j * 0.62)
        rect(s, x, iy, Inches(6.1), Inches(0.52), DARK)
        rect(s, x, iy + Inches(0.12), Inches(0.04), Inches(0.28), col)
        txbox(s, "☐  " + item, x + Inches(0.12), iy + Inches(0.07), Inches(5.8), Inches(0.4),
              size=10, color=WHITE)

y2 = Inches(2.1)
for title, col, items in col2:
    txbox(s, title, Inches(7.0), y2, Inches(6.0), Inches(0.3), size=9, bold=True, color=col)
    for j, item in enumerate(items):
        iy = y2 + Inches(0.45 + j * 0.62)
        rect(s, Inches(7.0), iy, Inches(6.0), Inches(0.52), DARK)
        rect(s, Inches(7.0), iy + Inches(0.12), Inches(0.04), Inches(0.28), col)
        txbox(s, "☐  " + item, Inches(7.12), iy + Inches(0.07), Inches(5.7), Inches(0.4),
              size=10, color=WHITE)
    y2 += Inches(0.45 + len(items) * 0.62 + 0.4)

footer_bar(s); slide_number(s, n)


prs.save(OUT)
print(f"✓ Saved: {OUT}")
print(f"  {n} slides · 13.33\" × 7.5\" widescreen")
