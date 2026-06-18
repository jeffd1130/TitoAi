#!/usr/bin/env python3
"""
Tito AI — Content Plan & Boost Budget Deck
W26–W28 · June–July 2026  (launch starts next week, Jun 22)
Progressive ramp: Week 1 TikTok → Week 2 TikTok+IG → Week 3 All 3
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT_PPTX = "/Users/jeff/Documents/Claude/TItoAi/content/strategy/TitoAI-ContentBoostPlan-W26-W28.pptx"
OUT_PDF  = "/Users/jeff/Documents/Claude/TItoAi/content/strategy/TitoAI-ContentBoostPlan-W26-W28.pdf"

# ── Brand colours ─────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0A, 0x0F, 0x1E)
GOLD   = RGBColor(0xF5, 0x9E, 0x0B)
TEAL   = RGBColor(0x0D, 0x94, 0x88)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x9C, 0xA3, 0xAF)
LGRAY  = RGBColor(0xE5, 0xE7, 0xEB)
DARK   = RGBColor(0x11, 0x18, 0x27)
DKDARK = RGBColor(0x0D, 0x16, 0x26)
RED    = RGBColor(0xDC, 0x26, 0x26)
GREEN  = RGBColor(0x06, 0x5F, 0x46)
LGREEN = RGBColor(0x10, 0xB9, 0x81)
PINK   = RGBColor(0xE8, 0x79, 0xF9)
BLUE   = RGBColor(0x93, 0xC5, 0xFD)
AMBER  = RGBColor(0xFC, 0xD3, 0x4D)
TIKRED = RGBColor(0xFF, 0x00, 0x50)
IGPUR  = RGBColor(0xC1, 0x3B, 0x84)
FBBLUE = RGBColor(0x18, 0x77, 0xF2)

# ── Layout constants ──────────────────────────────────────────────────────────
W     = Inches(13.33)
H     = Inches(7.5)
CTOP  = Inches(1.95)   # content top (after header)
CBOT  = Inches(6.85)   # content bottom (before footer)
CL    = Inches(0.3)    # left margin
CW    = Inches(12.73)  # content width
MID   = Inches(6.665)  # horizontal centre

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]
_sn   = [0]

# ── Primitive helpers ─────────────────────────────────────────────────────────
def add_slide():
    _sn[0] += 1
    return prs.slides.add_slide(BLANK)

def bg(s, color=NAVY):
    f = s.background.fill; f.solid(); f.fore_color.rgb = color

def box(s, l, t, w, h, fill, border=None):
    sh = s.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border; sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    return sh

def tx(s, text, l, t, w, h, sz=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = "Arial"
    return tb

def gold_bar(s):
    box(s, 0, 0, W, Pt(6), GOLD)

def footer(s):
    box(s, 0, H - Inches(0.44), W, Inches(0.44), DARK)
    tx(s, "Tito AI @TitoAIPH  ·  Content Plan & Boost Budget  ·  W26–W28 · June–July 2026",
       Inches(0.4), H - Inches(0.41), Inches(9), Inches(0.36), sz=7.5, color=GRAY)
    tx(s, str(_sn[0]),
       W - Inches(0.55), H - Inches(0.41), Inches(0.4), Inches(0.36),
       sz=8, color=GOLD, align=PP_ALIGN.RIGHT)

def hdr(s, title, sub=None, accent=GOLD):
    box(s, 0, Inches(1.06), Inches(0.07), Inches(0.54), accent)
    tx(s, title, Inches(0.22), Inches(1.04), Inches(12.8), Inches(0.66),
       sz=28, bold=True, color=WHITE)
    if sub:
        tx(s, sub, Inches(0.24), Inches(1.74), Inches(12.8), Inches(0.28),
           sz=9.5, italic=True, color=GRAY)
    box(s, Inches(0.22), Inches(1.83), Inches(12.9), Pt(2), accent)

def pill(s, text, x, y, w, h, bg_c, fg_c, sz=8.5):
    box(s, x, y, w, h, bg_c)
    tx(s, text, x + Inches(0.06), y + Inches(0.04),
       w - Inches(0.1), h - Inches(0.06), sz=sz, bold=True, color=fg_c,
       align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); gold_bar(s)
box(s, 0, H - Inches(2.2), W, Inches(2.2), DARK)

tx(s, "CONTENT PLAN & BOOST BUDGET · W26–W28 · JUNE–JULY 2026",
   Inches(0.6), Inches(0.42), Inches(12), Inches(0.38), sz=9, bold=True, color=GOLD)
tx(s, "Tito AI", Inches(0.6), Inches(0.9), Inches(12), Inches(1.0),
   sz=62, bold=True, color=WHITE)
tx(s, "@TitoAIPH", Inches(0.6), Inches(1.85), Inches(12), Inches(0.6),
   sz=36, bold=True, color=GOLD)
tx(s, "3 weeks · 9 posts · Php 6,000 total · Launch: Jun 22, 2026",
   Inches(0.6), Inches(2.56), Inches(12), Inches(0.44), sz=17, color=LGRAY)

box(s, Inches(0.6), Inches(3.14), Inches(0.7), Pt(4), GOLD)

# Ramp-up visual
for i, (wk, lbl, col) in enumerate([
    ("W26  Jun 22–28", "TikTok only",              TIKRED),
    ("W27  Jun 29–Jul 5", "TikTok + Instagram",     IGPUR),
    ("W28  Jul 6–12",  "All 3 platforms",           FBBLUE),
]):
    bx = Inches(0.6 + i * 4.1)
    box(s, bx, Inches(3.32), Inches(3.85), Inches(0.66), DARK)
    box(s, bx, Inches(3.32), Pt(4), Inches(0.66), col)
    tx(s, wk,  bx + Inches(0.12), Inches(3.38), Inches(2.0), Inches(0.28),
       sz=8.5, bold=True, color=col)
    tx(s, lbl, bx + Inches(0.12), Inches(3.58), Inches(3.6), Inches(0.28),
       sz=10, bold=True, color=WHITE)
    tx(s, "Php 2,000/week", bx + Inches(0.12), Inches(4.06),
       Inches(3.6), Inches(0.26), sz=10, bold=True, color=AMBER)

for i, (lbl, val) in enumerate([
    ("TOTAL BUDGET", "Php 6,000"),
    ("LAUNCH DATE",  "Jun 22, 2026"),
    ("GOAL",         "Followers + Views"),
]):
    bx = Inches(0.6 + i * 4.3)
    tx(s, lbl, bx, H - Inches(2.02), Inches(4.0), Inches(0.22), sz=7.5, bold=True, color=GRAY)
    tx(s, val,  bx, H - Inches(1.75), Inches(4.0), Inches(0.35), sz=14, color=WHITE)

footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — TARGET MARKET
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); gold_bar(s)
hdr(s, "Target Market — Mga Pamangkin", "Who we reach · Why they follow · Why they share")

# One-liner bar
box(s, CL, CTOP, CW, Inches(0.36), DKDARK)
box(s, CL, CTOP, Pt(4), Inches(0.36), GOLD)
tx(s, '"The warm, relatable Tito who teaches everyday Filipinos to use AI for free — before it replaces their job."',
   Inches(0.5), CTOP + Inches(0.05), CW - Inches(0.3), Inches(0.28),
   sz=9.5, italic=True, color=AMBER)

# Audience segments (left column — 6.5" wide)
SEG_Y = CTOP + Inches(0.45)
SEG_H = Inches(0.84)
SEG_GAP = Inches(0.06)
segments = [
    (GOLD,   "FREELANCERS",        "1.5M+ in PH",
     "Fear: clients will use AI instead. Need: faster proposals, better deliverables. High share rate."),
    (TEAL,   "GURO / TEACHERS",    "900K+ DepEd public school",
     "Fear: curriculum disruption. Need: lesson plans, grading, parent comms. One guro → 40 families."),
    (TIKRED, "BPO WORKERS",        "1.9M employees (BSP 2025)",
     "Fear: automation layoffs. Need: upskill NOW. Fastest-adopting audience for AI education content."),
    (PINK,   "NANAYS / TATAYS",    "Stay-at-home + working parents",
     "Fear: left behind. Need: family budgeting, homework help, meal planning, schedule management."),
    (FBBLUE, "SMALL BUSI. OWNERS", "SME operators — 14.9% use AI",
     "Fear: competition. Need: AI for marketing, customer service, operations. Highest spend intent."),
]
for i, (col, seg, stat, desc) in enumerate(segments):
    y = SEG_Y + i * (SEG_H + SEG_GAP)
    box(s, CL, y, Inches(6.5), SEG_H, DARK)
    box(s, CL, y, Pt(4), SEG_H, col)
    tx(s, seg,  Inches(0.5), y + Inches(0.08), Inches(2.4), Inches(0.28),
       sz=9.5, bold=True, color=col)
    tx(s, stat, Inches(2.95), y + Inches(0.08), Inches(3.8), Inches(0.26),
       sz=8.5, color=GRAY)
    tx(s, desc, Inches(0.5), y + Inches(0.38), Inches(6.0), Inches(0.4),
       sz=8.5, color=LGRAY)

# Right panel — psychographic profile
RX = Inches(7.1)
RW = Inches(5.93)

box(s, RX, CTOP + Inches(0.45), RW, Inches(0.3), DKDARK)
tx(s, "AUDIENCE PROFILE", RX + Inches(0.1), CTOP + Inches(0.5),
   RW - Inches(0.2), Inches(0.24), sz=8, bold=True, color=GOLD)

profile = [
    ("Age range",       "22–45 · Philippines · Taglish-speaking"),
    ("Pain point",      "Afraid AI will replace them before they learn to use it"),
    ("Motivation",      "Practical wins · Libre tools · No jargon · Real examples"),
    ("Share trigger",   "Posts that say 'this is for me' — specific personas, Php budgets, local context"),
    ("Watch trigger",   "Hooks that name their job or fear in the first 2 seconds"),
    ("CTA response",    "YES/NO questions · budget comments · 'I-try ko ito' replies"),
    ("Best platform",   "TikTok for discovery · Instagram for saves · Facebook for community"),
]
for i, (k, v) in enumerate(profile):
    y = CTOP + Inches(0.85 + i * 0.54)
    row_bg = DKDARK if i % 2 == 0 else DARK
    box(s, RX, y, RW, Inches(0.48), row_bg)
    tx(s, k, RX + Inches(0.1), y + Inches(0.08), Inches(1.55), Inches(0.3),
       sz=8, bold=True, color=TEAL)
    tx(s, v, RX + Inches(1.7), y + Inches(0.08), Inches(4.1), Inches(0.34),
       sz=8.5, color=LGRAY)

# Why they share bar
wsy = CTOP + Inches(0.85 + 7 * 0.54)
box(s, RX, wsy, RW, Inches(0.42), GREEN)
tx(s, "12.7M Filipinos exposed to AI disruption (ILO 2026) — almost none have a Taglish guide. Tito AI = first mover.",
   RX + Inches(0.1), wsy + Inches(0.07), RW - Inches(0.2), Inches(0.34),
   sz=8.5, bold=True, color=WHITE)

footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — STRATEGY OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); gold_bar(s)
hdr(s, "Boost Strategy Overview", "Php 2,000/week · 3 posts/week · progressive platform ramp-up · launch Jun 22")

# 5 overview rows
rows = [
    (GOLD,   "Weekly budget",      "Php 2,000/week · Php 6,000 total across 3 weeks (W26–W28)"),
    (TEAL,   "Posts per week",     "3 posts boosted each week — Mon AI Tip · Wed Demo · Fri Story"),
    (LGREEN, "Goal",               "Increase followers and video views. Target: 240–450 new followers over 3 weeks."),
    (PINK,   "Boost objective",    "Video views (TikTok Promote) · Profile visits (Instagram) · Reach + Page likes (Facebook)"),
    (AMBER,  "Audience targeting", "Philippines · Age 22–45 · Taglish · Tech / Freelancing / Education interests"),
]
for i, (col, k, v) in enumerate(rows):
    y = CTOP + i * Inches(0.60)
    box(s, CL, y, CW, Inches(0.52), DARK)
    box(s, CL, y, Pt(4), Inches(0.52), col)
    tx(s, k, Inches(0.5), y + Inches(0.07), Inches(2.7), Inches(0.24),
       sz=9.5, bold=True, color=col)
    tx(s, v, Inches(3.4), y + Inches(0.07), Inches(9.4), Inches(0.38),
       sz=9.5, color=LGRAY)

# Ramp visual — 5 rows × 0.60" = 3.0" → ends at 4.95"; VY=5.05, VH=1.78 → ends at 6.83"
VY = CTOP + Inches(3.10)   # = 5.05"
VH = Inches(1.78)           # ends at 6.83" ✓
box(s, CL, VY, CW, VH, DKDARK)
tx(s, "PLATFORM RAMP-UP  ·  One platform per week, add another each week until all 3 are running",
   Inches(0.5), VY + Inches(0.08), CW - Inches(0.4), Inches(0.22), sz=8.5, color=GRAY)

ramp = [
    ("WEEK 1",   "W26  Jun 22–28",  "TikTok only",
     "Learn TikTok Promote · baseline CPV + CPF · lowest risk", TIKRED),
    ("WEEK 2",   "W27  Jun 29–Jul 5","TikTok + Instagram",
     "Add IG Boost · compare cost-per-follow · two platforms", IGPUR),
    ("WEEK 3",   "W28  Jul 6–12",   "All 3 platforms",
     "Full push · FB Page required · scale winners in July", FBBLUE),
]
for i, (wk, dates, plats, desc, col) in enumerate(ramp):
    bx = Inches(0.5 + i * 4.26)
    by = VY + Inches(0.36)   # card starts 0.36" below ramp header
    bw = Inches(4.06)
    bh = Inches(1.36)        # 5.41 + 1.36 = 6.77" ✓ (under CBOT=6.85")
    box(s, bx, by, bw, bh, DARK)
    box(s, bx, by, bw, Pt(4), col)
    tx(s, wk,   bx + Inches(0.14), by + Inches(0.08), bw - Inches(0.2), Inches(0.24),
       sz=9, bold=True, color=col)
    tx(s, dates, bx + Inches(0.14), by + Inches(0.34), bw - Inches(0.2), Inches(0.22),
       sz=8, color=GRAY)
    box(s, bx + Inches(0.14), by + Inches(0.60), bw - Inches(0.3), Inches(0.28), DKDARK)
    tx(s, plats, bx + Inches(0.2), by + Inches(0.63), bw - Inches(0.4), Inches(0.24),
       sz=9.5, bold=True, color=col, align=PP_ALIGN.CENTER)
    tx(s, desc, bx + Inches(0.14), by + Inches(0.94), bw - Inches(0.2), Inches(0.32),
       sz=8.5, color=LGRAY)
    tx(s, "Php 2,000", bx + bw - Inches(1.6), by + bh - Inches(0.30),
       Inches(1.5), Inches(0.24), sz=9, bold=True, color=AMBER, align=PP_ALIGN.RIGHT)

footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — 9-POST CONTENT CALENDAR
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); gold_bar(s)
hdr(s, "3-Week Content Calendar", "9 posts · Mon AI Tip / Wed Demo / Fri Story · drop 7–8 PM PHT · launch Jun 22")

# Column widths
CWS = [Inches(0.78), Inches(1.15), Inches(1.12), Inches(3.88), Inches(1.72), Inches(1.44), Inches(2.64)]
CXS = [CL]
for w in CWS[:-1]: CXS.append(CXS[-1] + w)

HDR_H  = Inches(0.34)
ROW_H  = Inches(0.46)
TOT_H  = Inches(0.32)
TABLE_H = HDR_H + 9 * ROW_H + TOT_H   # 0.34 + 4.14 + 0.32 = 4.80 ✓

# Header row
for j, (h, x, w) in enumerate(zip(
    ["Week", "Drop Date", "Day", "Post / Topic", "Platform", "Budget", "Boost Dates"],
    CXS, CWS
)):
    box(s, x, CTOP, w, HDR_H, GREEN)
    tx(s, h, x + Inches(0.05), CTOP + Inches(0.06), w - Inches(0.08), Inches(0.24),
       sz=8, bold=True, color=WHITE)

# Posts
posts = [
    ("W26", TIKRED, "Jun 22","MON","Ang Sabi Nila: Pang-Matalino Lang Iyan",
     "TikTok", TIKRED,"Php 700","Jun 23–24"),
    ("W26", TIKRED, "Jun 25","WED","Gemini para sa Guro",
     "TikTok", TIKRED,"Php 700","Jun 26–27"),
    ("W26", TIKRED, "Jun 27","FRI","1 Buwan Kasama Tito AI — Ang Resulta",
     "TikTok", TIKRED,"Php 600","Jun 28–29"),
    ("W27", IGPUR,  "Jun 30","MON","July Teaser — Ano ang Susunod?",
     "TikTok", TIKRED,"Php 700","Jul 1–2"),
    ("W27", IGPUR,  "Jul 2", "WED","TBD — Wednesday Demo",
     "Instagram",IGPUR,"Php 700","Jul 3–4"),
    ("W27", IGPUR,  "Jul 4", "FRI","TBD — Friday Story / Inspiration",
     "TikTok", TIKRED,"Php 600","Jul 5–6"),
    ("W28", FBBLUE, "Jul 7", "MON","TBD — Monday AI Tip",
     "TikTok", TIKRED,"Php 600","Jul 8–9"),
    ("W28", FBBLUE, "Jul 9", "WED","TBD — Wednesday Demo",
     "Facebook",FBBLUE,"Php 800","Jul 10–11"),
    ("W28", FBBLUE, "Jul 11","FRI","TBD — Friday Story",
     "Instagram",IGPUR,"Php 600","Jul 12–13"),
]
for i, (wk,wc, date, day, topic, plat, pc, bgt, bdates) in enumerate(posts):
    y = CTOP + HDR_H + i * ROW_H
    rb = DKDARK if i % 2 == 0 else DARK
    for x, w in zip(CXS, CWS):
        box(s, x, y, w, ROW_H, rb)
    # Week pill
    box(s, CXS[0], y, CWS[0], ROW_H, wc)
    tx(s, wk, CXS[0]+Inches(0.04), y+Inches(0.1), CWS[0]-Inches(0.06), Inches(0.28),
       sz=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Date
    tx(s, date, CXS[1]+Inches(0.05), y+Inches(0.1), CWS[1]-Inches(0.08), Inches(0.28),
       sz=8.5, bold=True, color=GOLD)
    # Day
    tx(s, day, CXS[2]+Inches(0.05), y+Inches(0.1), CWS[2]-Inches(0.08), Inches(0.28),
       sz=8, color=GRAY)
    # Topic
    tx(s, topic, CXS[3]+Inches(0.05), y+Inches(0.06), CWS[3]-Inches(0.08), Inches(0.36),
       sz=8.5, bold=True, color=WHITE)
    # Platform pill
    box(s, CXS[4]+Inches(0.06), y+Inches(0.09), CWS[4]-Inches(0.12), Inches(0.28), pc)
    tx(s, plat, CXS[4]+Inches(0.06), y+Inches(0.11), CWS[4]-Inches(0.12), Inches(0.26),
       sz=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Budget
    tx(s, bgt, CXS[5]+Inches(0.05), y+Inches(0.1), CWS[5]-Inches(0.08), Inches(0.28),
       sz=9, bold=True, color=LGREEN)
    # Boost dates
    tx(s, bdates, CXS[6]+Inches(0.05), y+Inches(0.1), CWS[6]-Inches(0.08), Inches(0.28),
       sz=8.5, color=LGRAY)

# Total row
ty = CTOP + HDR_H + 9 * ROW_H
for x, w in zip(CXS, CWS):
    box(s, x, ty, w, TOT_H, GREEN)
tx(s, "TOTAL", CXS[0]+Inches(0.05), ty+Inches(0.05), Inches(3.0), Inches(0.22),
   sz=8.5, bold=True, color=WHITE)
tx(s, "9 posts", CXS[3]+Inches(0.05), ty+Inches(0.05), Inches(2.0), Inches(0.22),
   sz=8.5, color=WHITE)
tx(s, "Php 6,000", CXS[5]+Inches(0.05), ty+Inches(0.04), CWS[5]-Inches(0.08), Inches(0.26),
   sz=10, bold=True, color=AMBER)
tx(s, "3 weeks · 3 platforms", CXS[6]+Inches(0.05), ty+Inches(0.05),
   CWS[6]-Inches(0.08), Inches(0.22), sz=7.5, color=WHITE)

footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — WEEK 1: W26 TIKTOK ONLY
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); gold_bar(s)
box(s, 0, Inches(1.06), Inches(0.07), Inches(0.54), TIKRED)
tx(s, "Week 1 — W26 · Jun 22–28 · TikTok Only", Inches(0.22), Inches(1.04),
   Inches(12.8), Inches(0.66), sz=27, bold=True, color=WHITE)
tx(s, "Learn TikTok Promote mechanics · establish baseline cost-per-view and cost-per-follower",
   Inches(0.24), Inches(1.74), Inches(12.8), Inches(0.28), sz=9.5, italic=True, color=GRAY)
box(s, Inches(0.22), Inches(1.83), Inches(12.9), Pt(2), TIKRED)

# ── Left: Budget breakdown ──
LW = Inches(6.1)
BHD_H = Inches(0.3)
POST_H = Inches(0.78)
TOT_H2 = Inches(0.32)

box(s, CL, CTOP, LW, BHD_H, TIKRED)
tx(s, "🎵  TIKTOK PROMOTE — ALL 3 POSTS  ·  Php 2,000 total",
   Inches(0.5), CTOP + Inches(0.05), LW - Inches(0.3), Inches(0.22),
   sz=9, bold=True, color=WHITE)

w1_posts = [
    ("MON · Jun 22", "Ang Sabi Nila: Pang-Matalino Lang Iyan",
     "Php 350/day × 2 days", "Php 700", "Jun 23–24"),
    ("WED · Jun 25", "Gemini para sa Guro",
     "Php 350/day × 2 days", "Php 700", "Jun 26–27"),
    ("FRI · Jun 27", "1 Buwan Kasama Tito AI — Ang Resulta",
     "Php 300/day × 2 days", "Php 600", "Jun 28–29"),
]
for i, (day, topic, rate, total, bdates) in enumerate(w1_posts):
    y = CTOP + BHD_H + i * POST_H
    box(s, CL, y, LW, POST_H, DKDARK if i%2==0 else DARK)
    box(s, CL, y, Pt(4), POST_H, TIKRED)
    tx(s, day,   Inches(0.5), y+Inches(0.08), Inches(1.6), Inches(0.26),
       sz=9, bold=True, color=TIKRED)
    tx(s, topic, Inches(0.5), y+Inches(0.36), Inches(4.0), Inches(0.32),
       sz=9, color=WHITE)
    tx(s, rate,  Inches(2.3), y+Inches(0.08), Inches(2.0), Inches(0.24),
       sz=8, color=GRAY)
    tx(s, f"Boost: {bdates}", Inches(2.3), y+Inches(0.36), Inches(2.0), Inches(0.28),
       sz=8, color=GRAY)
    box(s, Inches(4.55), y+Inches(0.14), Inches(1.6), Inches(0.5), DKDARK if i%2==1 else DARK)
    tx(s, total, Inches(4.55), y+Inches(0.18), Inches(1.6), Inches(0.36),
       sz=14, bold=True, color=LGREEN, align=PP_ALIGN.CENTER)

# Total
ty1 = CTOP + BHD_H + 3 * POST_H
box(s, CL, ty1, LW, TOT_H2, GREEN)
tx(s, "W26 WEEK TOTAL:  Php 2,000", Inches(0.5), ty1+Inches(0.06),
   LW - Inches(0.3), Inches(0.24), sz=10, bold=True, color=WHITE)

# ── Right: Setup guide ──
RX2 = Inches(6.7)
RW2 = Inches(6.33)
SETUP_H = BHD_H + 3 * POST_H + TOT_H2  # same height as left block

box(s, RX2, CTOP, RW2, BHD_H, DKDARK)
tx(s, "HOW TO SET UP — TikTok Promote",
   RX2 + Inches(0.1), CTOP + Inches(0.06), RW2 - Inches(0.2), Inches(0.22),
   sz=9, bold=True, color=TIKRED)

steps = [
    ("1", "Open TikTok app → go to the posted video"),
    ("2", "Tap share icon (arrow) → tap \"Promote\""),
    ("3", "Goal: select \"More video views\" — NOT followers. Views are cheaper and drive the algorithm, which brings followers organically."),
    ("4", "Audience: select \"Automatic\" — Smart targeting outperforms manual on new accounts"),
    ("5", "Duration: 2 days · Budget: Php 300–350/day · Confirm payment method"),
    ("6", "Start: morning AFTER posting (let organic play 8–12 hrs first for best CPV)"),
]
step_h = (SETUP_H - BHD_H) / 6
for i, (num, step) in enumerate(steps):
    y = CTOP + BHD_H + i * step_h
    box(s, RX2, y, RW2, step_h, DKDARK if i%2==0 else DARK)
    box(s, RX2, y, Inches(0.34), step_h, TIKRED)
    tx(s, num, RX2+Inches(0.04), y+Inches(0.05), Inches(0.28), step_h-Inches(0.06),
       sz=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(s, step, RX2+Inches(0.42), y+Inches(0.06), RW2-Inches(0.52), step_h-Inches(0.1),
       sz=8.5, color=LGRAY)

# Bottom KPI bar
kpi_y = CTOP + SETUP_H + Inches(0.1)
box(s, CL, kpi_y, CW, Inches(0.34), DKDARK)
box(s, CL, kpi_y, Pt(4), Inches(0.34), GOLD)
tx(s, "W26 TARGETS:  Reach 9,000+  ·  Views 15,000+  ·  New followers 60–120  ·  Cost-per-follow target < Php 25",
   Inches(0.5), kpi_y+Inches(0.07), CW-Inches(0.3), Inches(0.24),
   sz=9, bold=True, color=AMBER)

# Tips
tip_y = kpi_y + Inches(0.42)
box(s, CL, tip_y, CW, Inches(0.48), DARK)
tips = "💡  Tips: Use a strong hook in the first 2 seconds · film face-cam for highest CTR · end with a YES/NO question to boost comment rate · reply to all comments in the first hour after posting"
tx(s, tips, Inches(0.5), tip_y+Inches(0.08), CW-Inches(0.3), Inches(0.36),
   sz=8.5, italic=True, color=LGRAY)

footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — WEEK 2: W27 TIKTOK + INSTAGRAM
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); gold_bar(s)
box(s, 0, Inches(1.06), Inches(0.07), Inches(0.54), IGPUR)
tx(s, "Week 2 — W27 · Jun 29–Jul 5 · TikTok + Instagram", Inches(0.22), Inches(1.04),
   Inches(12.8), Inches(0.66), sz=27, bold=True, color=WHITE)
tx(s, "Apply W26 learnings · add Instagram Boost · compare cost-per-follow across two platforms",
   Inches(0.24), Inches(1.74), Inches(12.8), Inches(0.28), sz=9.5, italic=True, color=GRAY)
box(s, Inches(0.22), Inches(1.83), Inches(12.9), Pt(2), IGPUR)

# Budget breakdown — left column
LW2 = Inches(6.2)
BHD = Inches(0.3)
PR_H = Inches(0.72)

box(s, CL, CTOP, LW2, BHD, DKDARK)
tx(s, "BUDGET BREAKDOWN — 3 POSTS",
   Inches(0.5), CTOP+Inches(0.06), LW2-Inches(0.3), Inches(0.22),
   sz=9, bold=True, color=IGPUR)

w2_posts = [
    ("MON · Jun 30", "July Teaser — Ano ang Susunod?",
     "TikTok Promote", TIKRED, "Php 350/day × 2 days", "Php 700", "Jul 1–2"),
    ("WED · Jul 2", "TBD — Wednesday Demo",
     "Instagram Boost", IGPUR, "Php 350/day × 2 days", "Php 700", "Jul 3–4"),
    ("FRI · Jul 4", "TBD — Friday Story / Inspiration",
     "TikTok Promote", TIKRED, "Php 300/day × 2 days", "Php 600", "Jul 5–6"),
]
for i, (day, topic, plat, pc, rate, total, bdates) in enumerate(w2_posts):
    y = CTOP + BHD + i * PR_H
    box(s, CL, y, LW2, PR_H, DKDARK if i%2==0 else DARK)
    box(s, CL, y, Pt(4), PR_H, pc)
    tx(s, day,   Inches(0.5), y+Inches(0.07), Inches(1.6), Inches(0.24),
       sz=8.5, bold=True, color=pc)
    tx(s, topic, Inches(0.5), y+Inches(0.34), Inches(3.6), Inches(0.3),
       sz=8.5, color=WHITE)
    box(s, Inches(2.2), y+Inches(0.07), Inches(1.55), Inches(0.24), pc)
    tx(s, plat, Inches(2.2), y+Inches(0.08), Inches(1.55), Inches(0.22),
       sz=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(s, rate, Inches(3.82), y+Inches(0.07), Inches(1.5), Inches(0.24),
       sz=7.5, color=GRAY)
    tx(s, f"Boost: {bdates}", Inches(3.82), y+Inches(0.36), Inches(1.5), Inches(0.26),
       sz=7.5, color=GRAY)
    box(s, Inches(5.25), y+Inches(0.12), Inches(0.8), Inches(0.46), DKDARK if i%2==1 else DARK)
    tx(s, total, Inches(5.25), y+Inches(0.16), Inches(0.8), Inches(0.36),
       sz=12, bold=True, color=LGREEN, align=PP_ALIGN.CENTER)

ty2 = CTOP + BHD + 3 * PR_H
box(s, CL, ty2, LW2, Inches(0.32), GREEN)
tx(s, "W27 WEEK TOTAL:  Php 2,000  (TikTok Php 1,300  +  Instagram Php 700)",
   Inches(0.5), ty2+Inches(0.07), LW2-Inches(0.3), Inches(0.24),
   sz=9, bold=True, color=WHITE)

# Right: Setup guides
RX3 = Inches(6.8)
RW3 = Inches(6.23)

# TikTok setup (compact)
tk_y = CTOP
box(s, RX3, tk_y, RW3, BHD, TIKRED)
tx(s, "🎵 TikTok Promote — same steps as Week 1",
   RX3+Inches(0.1), tk_y+Inches(0.06), RW3-Inches(0.2), Inches(0.22),
   sz=8.5, bold=True, color=WHITE)
tik_notes = [
    ("Goal", "More video views · Duration: 2 days · Php 300–350/day"),
    ("Audience", "Automatic · Start boost morning after posting"),
    ("New in W27", "Check W26 results first: if CPF < Php 20 on a video, increase to Php 400/day"),
]
for i, (k, v) in enumerate(tik_notes):
    y = tk_y + BHD + i * Inches(0.44)
    box(s, RX3, y, RW3, Inches(0.42), DKDARK if i%2==0 else DARK)
    tx(s, k, RX3+Inches(0.1), y+Inches(0.08), Inches(1.3), Inches(0.26),
       sz=8, bold=True, color=TIKRED)
    tx(s, v, RX3+Inches(1.45), y+Inches(0.08), RW3-Inches(1.55), Inches(0.3),
       sz=8.5, color=LGRAY)

# Instagram Boost setup
ig_y = tk_y + BHD + 3 * Inches(0.44) + Inches(0.12)
box(s, RX3, ig_y, RW3, BHD, IGPUR)
tx(s, "📸 Instagram Boost — NEW in Week 2",
   RX3+Inches(0.1), ig_y+Inches(0.06), RW3-Inches(0.2), Inches(0.22),
   sz=8.5, bold=True, color=WHITE)
ig_steps = [
    ("1", "Open Instagram → go to the posted Reel/carousel → tap \"Boost Post\""),
    ("2", "Goal: \"More profile visits\" — drives follow intent (not just views)"),
    ("3", "Audience: Automatic · Philippines auto-detected"),
    ("4", "Duration: 2 days · Budget: Php 350/day · Confirm payment"),
    ("5", "Check results in Instagram Insights → Boosted post tab"),
]
for i, (num, step) in enumerate(ig_steps):
    y = ig_y + BHD + i * Inches(0.44)
    box(s, RX3, y, RW3, Inches(0.42), DKDARK if i%2==0 else DARK)
    box(s, RX3, y, Inches(0.32), Inches(0.42), IGPUR)
    tx(s, num, RX3+Inches(0.04), y+Inches(0.1), Inches(0.26), Inches(0.26),
       sz=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(s, step, RX3+Inches(0.4), y+Inches(0.07), RW3-Inches(0.5), Inches(0.3),
       sz=8.5, color=LGRAY)

# KPI bar
kpi2_y = max(ty2, ig_y + BHD + 5 * Inches(0.44)) + Inches(0.12)
box(s, CL, kpi2_y, CW, Inches(0.32), DKDARK)
box(s, CL, kpi2_y, Pt(4), Inches(0.32), GOLD)
tx(s, "W27 TARGETS:  TikTok 10,000+ views  ·  IG 3,000+ reach  ·  Combined new followers 80–160  ·  Compare CPF: which platform wins?",
   Inches(0.5), kpi2_y+Inches(0.07), CW-Inches(0.3), Inches(0.24),
   sz=9, bold=True, color=AMBER)

footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — WEEK 3: W28 ALL 3 PLATFORMS
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); gold_bar(s)
box(s, 0, Inches(1.06), Inches(0.07), Inches(0.54), FBBLUE)
tx(s, "Week 3 — W28 · Jul 6–12 · All 3 Platforms", Inches(0.22), Inches(1.04),
   Inches(12.8), Inches(0.66), sz=27, bold=True, color=WHITE)
tx(s, "Full 3-platform push · compare all results · decide where to scale in July · requires Facebook Page",
   Inches(0.24), Inches(1.74), Inches(12.8), Inches(0.28), sz=9.5, italic=True, color=GRAY)
box(s, Inches(0.22), Inches(1.83), Inches(12.9), Pt(2), FBBLUE)

# 3-column card layout
CARD_W = Inches(4.12)
CARD_H = Inches(3.38)
CARD_GAP = Inches(0.24)

platforms_w3 = [
    ("🎵 TikTok Promote", TIKRED, "MON · Jul 7", "TBD — Monday AI Tip",
     "Php 300/day × 2 days", "Php 600", "Jul 8–9",
     "Video views · Automatic audience\nStart Jul 8 morning"),
    ("📸 Instagram Boost", IGPUR, "FRI · Jul 11", "TBD — Friday Story",
     "Php 300/day × 2 days", "Php 600", "Jul 12–13",
     "Profile visits · Automatic audience\nStart Jul 12 morning"),
    ("📘 Facebook Boost", FBBLUE, "WED · Jul 9", "TBD — Wednesday Demo",
     "Php 400/day × 2 days", "Php 800", "Jul 10–11",
     "Reach + Page likes · PH 22–45\nTech + Freelancing interests"),
]
for i, (plat, pc, day, topic, rate, total, bdates, obj) in enumerate(platforms_w3):
    cx = CL + i * (CARD_W + CARD_GAP)
    cy = CTOP
    box(s, cx, cy, CARD_W, CARD_H, DARK)
    box(s, cx, cy, CARD_W, Pt(4), pc)
    tx(s, plat,  cx+Inches(0.12), cy+Inches(0.1), CARD_W-Inches(0.2), Inches(0.28),
       sz=9.5, bold=True, color=pc)
    # Budget badge
    box(s, cx+Inches(0.12), cy+Inches(0.44), CARD_W-Inches(0.24), Inches(0.38), DKDARK)
    tx(s, total, cx+Inches(0.12), cy+Inches(0.5), CARD_W-Inches(0.24), Inches(0.3),
       sz=16, bold=True, color=LGREEN, align=PP_ALIGN.CENTER)
    tx(s, day,   cx+Inches(0.12), cy+Inches(0.96), CARD_W-Inches(0.2), Inches(0.26),
       sz=9, bold=True, color=pc)
    tx(s, topic, cx+Inches(0.12), cy+Inches(1.26), CARD_W-Inches(0.2), Inches(0.4),
       sz=9, color=WHITE)
    tx(s, rate,  cx+Inches(0.12), cy+Inches(1.72), CARD_W-Inches(0.2), Inches(0.26),
       sz=8.5, color=GOLD)
    tx(s, f"Boost: {bdates}", cx+Inches(0.12), cy+Inches(2.0), CARD_W-Inches(0.2), Inches(0.26),
       sz=8.5, color=GRAY)
    box(s, cx+Inches(0.12), cy+Inches(2.32), CARD_W-Inches(0.24), Pt(1), DKDARK)
    tx(s, obj, cx+Inches(0.12), cy+Inches(2.42), CARD_W-Inches(0.2), Inches(0.7),
       sz=8.5, italic=True, color=LGRAY)

# Facebook Page warning
warn_y = CTOP + CARD_H + Inches(0.12)
box(s, CL, warn_y, CW, Inches(0.56), RGBColor(0x7F, 0x1D, 0x1D))
box(s, CL, warn_y, Pt(5), Inches(0.56), RED)
tx(s, "⚠️  BEFORE WEEK 3: Facebook Page required. Go to facebook.com/pages/create → \"Business or brand\" → migrate followers. Personal profiles CANNOT boost or access analytics.",
   Inches(0.5), warn_y+Inches(0.1), CW-Inches(0.3), Inches(0.4),
   sz=9, bold=True, color=WHITE)

# Grand total
tot_y = warn_y + Inches(0.64)
box(s, CL, tot_y, CW, Inches(0.38), GREEN)
tx(s, "W28 TOTAL:  Php 2,000  ·  TikTok Php 600  +  Instagram Php 600  +  Facebook Php 800  ·  3-Week Grand Total: Php 6,000",
   Inches(0.5), tot_y+Inches(0.08), CW-Inches(0.3), Inches(0.26),
   sz=9.5, bold=True, color=WHITE)

# Scale note
sc_y = tot_y + Inches(0.46)
box(s, CL, sc_y, CW, Inches(0.38), DKDARK)
box(s, CL, sc_y, Pt(4), Inches(0.38), GOLD)
tx(s, "After W28: compare CPF across platforms → winner gets biggest July budget. Scale rule: CPF < Php 12 → add Php 400/week to that platform.",
   Inches(0.5), sc_y+Inches(0.08), CW-Inches(0.3), Inches(0.26),
   sz=8.5, color=AMBER)

footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — BUDGET SUMMARY
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); gold_bar(s)
hdr(s, "Budget Summary — 3 Weeks", "Php 6,000 total · Php 2,000/week · W26 TikTok → W27 TikTok+IG → W28 All 3")

# Table columns
TC_W = [Inches(1.22), Inches(2.54), Inches(2.54), Inches(2.18), Inches(1.56), Inches(1.74), Inches(0.95)]
TC_X = [CL]
for w in TC_W[:-1]: TC_X.append(TC_X[-1]+w)
HDR_ROW_H = Inches(0.36)
DATA_ROW_H = Inches(0.82)
TOT_ROW_H  = Inches(0.36)

# Header
for j, (h, x, w) in enumerate(zip(
    ["Week", "TikTok Spend", "Instagram Spend", "Facebook Spend", "Weekly Total", "Followers Target", "Platforms"],
    TC_X, TC_W
)):
    box(s, x, CTOP, w, HDR_ROW_H, DARK)
    tx(s, h, x+Inches(0.06), CTOP+Inches(0.07), w-Inches(0.1), Inches(0.24),
       sz=8, bold=True, color=GOLD)

# Data rows
budget_rows = [
    ("W26\nJun 22–28",
     "Php 2,000\n3 posts × TikTok",
     "—", "—",
     "Php 2,000", "60–120", "1"),
    ("W27\nJun 29–Jul 5",
     "Php 1,300\n2 posts × TikTok",
     "Php 700\n1 post × Instagram",
     "—",
     "Php 2,000", "80–160", "2"),
    ("W28\nJul 6–12",
     "Php 600\n1 post × TikTok",
     "Php 600\n1 post × Instagram",
     "Php 800\n1 post × Facebook",
     "Php 2,000", "100–180", "3"),
]
VCOLS = [None, TIKRED, IGPUR, FBBLUE, LGREEN, TEAL, GOLD]
for i, row in enumerate(budget_rows):
    y = CTOP + HDR_ROW_H + i * DATA_ROW_H
    rb = DKDARK if i%2==0 else DARK
    for j, (val, x, w, vc) in enumerate(zip(row, TC_X, TC_W, VCOLS)):
        box(s, x, y, w, DATA_ROW_H, rb)
        c = vc if vc and val != "—" else (GRAY if val == "—" else WHITE)
        if j == 0: c = [TIKRED, IGPUR, FBBLUE][i]
        tx(s, val, x+Inches(0.07), y+Inches(0.14), w-Inches(0.1), DATA_ROW_H-Inches(0.2),
           sz=9, bold=(j in [0, 4]), color=c)

# Grand total
gt_y = CTOP + HDR_ROW_H + 3 * DATA_ROW_H
box(s, CL, gt_y, CW, TOT_ROW_H, GREEN)
totals = ["TOTAL", "Php 3,900", "Php 1,300", "Php 800", "Php 6,000", "240–460", "—"]
tcols2 = [WHITE, TIKRED, IGPUR, FBBLUE, AMBER, LGREEN, WHITE]
for val, x, w, c in zip(totals, TC_X, TC_W, tcols2):
    tx(s, val, x+Inches(0.07), gt_y+Inches(0.07), w-Inches(0.1), Inches(0.24),
       sz=9, bold=True, color=c)

# Scale rules
sr_y = gt_y + TOT_ROW_H + Inches(0.16)
box(s, CL, sr_y, CW, Inches(0.3), DKDARK)
box(s, CL, sr_y, Pt(4), Inches(0.3), GOLD)
tx(s, "SCALE RULES — WEEK 4+ (July)", Inches(0.5), sr_y+Inches(0.06),
   Inches(4.0), Inches(0.22), sz=8.5, bold=True, color=GOLD)

scale_rules = [
    ("Single post earns 50+ followers", "→ Extend that post 2 more days at same daily rate"),
    ("Cost-per-follow drops below Php 12", "→ Increase that platform's weekly budget by Php 400"),
    ("Platform with lowest CPF", "→ Gets biggest July budget allocation (min Php 800/week)"),
    ("Post earns < 20 followers after boost", "→ Swap to a different platform or content slot next week"),
]
sr_row_h = Inches(0.36)
for i, (trigger, action) in enumerate(scale_rules):
    y = sr_y + Inches(0.3) + i * sr_row_h
    box(s, CL, y, CW, sr_row_h, DKDARK if i%2==0 else DARK)
    tx(s, trigger, Inches(0.5), y+Inches(0.07), Inches(5.5), Inches(0.26),
       sz=8.5, color=LGRAY)
    tx(s, action,  Inches(6.1), y+Inches(0.07), Inches(6.8), Inches(0.26),
       sz=8.5, bold=True, color=AMBER)

footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — KPIs & TRACKING
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); gold_bar(s)
hdr(s, "KPIs & Weekly Tracking", "Pull every Sunday · share screenshots with Jeff · review before next week's boost")

# 3 platform blocks
BLK_W = Inches(4.12)
BLK_GAP = Inches(0.24)
BLK_HDR_H = Inches(0.3)
BLK_SUB_H = Inches(0.24)
KPI_ROW_H = Inches(0.42)
N_KPIS = 5
BLK_H = BLK_HDR_H + BLK_SUB_H + N_KPIS * KPI_ROW_H  # 0.3+0.24+2.1 = 2.64"

plat_kpis = [
    ("TikTok", TIKRED, [
        ("Reach per boosted post", "3,000+", "7,000+"),
        ("Video views per post",   "5,000+", "15,000+"),
        ("New followers (3 posts)","60–120", "250+"),
        ("Cost per follow",        "< Php 25", "< Php 15"),
        ("Watch completion rate",  "> 35%",  "> 50%"),
    ]),
    ("Instagram", IGPUR, [
        ("Reach per boosted post", "2,000+", "5,000+"),
        ("Profile visits per post","200+",   "600+"),
        ("New followers (per post)","20–40", "100+"),
        ("Cost per follow",        "< Php 25", "< Php 15"),
        ("Saves per post",         "20+",    "80+"),
    ]),
    ("Facebook", FBBLUE, [
        ("Reach per boosted post", "2,500+", "6,000+"),
        ("New Page likes per post","20–40",  "100+"),
        ("Cost per Page like",     "< Php 25", "< Php 15"),
        ("Link clicks to profile", "50+",   "200+"),
        ("Active weeks",           "W28 only", "Full July"),
    ]),
]

for i, (plat, pc, kpis) in enumerate(plat_kpis):
    bx = CL + i * (BLK_W + BLK_GAP)
    # Header
    box(s, bx, CTOP, BLK_W, BLK_HDR_H, pc)
    tx(s, plat, bx+Inches(0.1), CTOP+Inches(0.06), BLK_W-Inches(0.2), Inches(0.22),
       sz=10, bold=True, color=WHITE)
    # Sub-header
    box(s, bx, CTOP+BLK_HDR_H, BLK_W, BLK_SUB_H, DKDARK)
    for j, col_lbl in enumerate(["Metric", "W1 Target", "W3 Target"]):
        cx2 = bx + [Inches(0.08), Inches(2.2), Inches(3.3)][j]
        cw2 = [Inches(2.1), Inches(1.0), Inches(0.9)][j]
        tx(s, col_lbl, cx2, CTOP+BLK_HDR_H+Inches(0.04), cw2, Inches(0.18),
           sz=7, bold=True, color=GRAY)
    # KPI rows
    for k, (metric, t1, t3) in enumerate(kpis):
        ky = CTOP + BLK_HDR_H + BLK_SUB_H + k * KPI_ROW_H
        row_bg = DKDARK if k%2==0 else DARK
        box(s, bx, ky, BLK_W, KPI_ROW_H, row_bg)
        tx(s, metric, bx+Inches(0.08), ky+Inches(0.08), Inches(2.1), Inches(0.28),
           sz=8, color=LGRAY)
        tx(s, t1, bx+Inches(2.2), ky+Inches(0.08), Inches(1.0), Inches(0.28),
           sz=8.5, bold=True, color=LGREEN, align=PP_ALIGN.CENTER)
        tx(s, t3, bx+Inches(3.3), ky+Inches(0.08), Inches(0.9), Inches(0.28),
           sz=8.5, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Where to pull
wp_y = CTOP + BLK_H + Inches(0.14)
box(s, CL, wp_y, CW, Inches(0.34), DARK)
box(s, CL, wp_y, Pt(4), Inches(0.34), TEAL)
sources = [
    ("TikTok", "Creator Center → Analytics + Promote tab", TIKRED),
    ("Instagram", "Professional Dashboard → Boosted post results", IGPUR),
    ("Facebook", "Business Suite → Page Insights (available once Page is created)", FBBLUE),
]
tx(s, "WHERE TO PULL:  ", Inches(0.5), wp_y+Inches(0.09), Inches(1.8), Inches(0.22),
   sz=8, bold=True, color=TEAL)
pull_str = "  ·  ".join([f"{p}: {src}" for p, src, _ in sources])
tx(s, pull_str, Inches(2.1), wp_y+Inches(0.09), CW-Inches(1.9), Inches(0.22),
   sz=8, color=LGRAY)

# Engagement checklist
ec_y = wp_y + Inches(0.42)
box(s, CL, ec_y, CW, Inches(0.28), DKDARK)
box(s, CL, ec_y, Pt(4), Inches(0.28), GOLD)
tx(s, "POST-BOOST ENGAGEMENT CHECKLIST (every post day)",
   Inches(0.5), ec_y+Inches(0.06), Inches(6.0), Inches(0.2), sz=8, bold=True, color=GOLD)

checklist = [
    "Reply to ALL comments within first 1 hour of posting",
    "Pin first comment within 5 min (tip or claude.ai link)",
    "Share Reel to Instagram Stories immediately",
    "Post to Facebook cross-post (native upload — never TikTok link)",
    "End every video with a YES/NO or choice question to drive comment rate",
]
CHK_H = Inches(0.27)
for i, item in enumerate(checklist):
    y = ec_y + Inches(0.28) + i * CHK_H
    box(s, CL + i%5 * Inches(0.0), y, CW, CHK_H, DKDARK if i%2==0 else DARK)
    tx(s, f"☐  {item}", Inches(0.5), y+Inches(0.07), CW-Inches(0.3), Inches(0.24),
       sz=8.5, color=LGRAY)

footer(s)


# ── Save PPTX ─────────────────────────────────────────────────────────────────
prs.save(OUT_PPTX)
print(f"✓ PPTX  {_sn[0]} slides → {OUT_PPTX}")


# ═══════════════════════════════════════════════════════════════════════════════
# PDF via ReportLab
# ═══════════════════════════════════════════════════════════════════════════════
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm
from reportlab.lib import colors
from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
                                TableStyle, HRFlowable, KeepTogether)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT

C_NAVY = colors.HexColor("#0A0F1E")
C_GOLD = colors.HexColor("#F59E0B")
C_TEAL = colors.HexColor("#0D9488")
C_WHITE= colors.white
C_GRAY = colors.HexColor("#9CA3AF")
C_LG   = colors.HexColor("#E5E7EB")
C_DARK = colors.HexColor("#111827")
C_DK2  = colors.HexColor("#0D1626")
C_RED  = colors.HexColor("#DC2626")
C_GRN  = colors.HexColor("#059669")
C_LGRN = colors.HexColor("#10B981")
C_PINK = colors.HexColor("#E879F9")
C_BLUE = colors.HexColor("#93C5FD")
C_AMBE = colors.HexColor("#FCD34D")
C_TIK  = colors.HexColor("#FF0050")
C_IGP  = colors.HexColor("#C13B84")
C_FBB  = colors.HexColor("#1877F2")

styles = getSampleStyleSheet()
def S(name, **kw):
    return ParagraphStyle(name+str(id(kw)), parent=styles.get(name, styles["Normal"]), **kw)

H2   = S("Normal", fontSize=14, leading=18, fontName="Helvetica-Bold", textColor=C_GOLD)
H3   = S("Normal", fontSize=10, leading=13, fontName="Helvetica-Bold", textColor=C_GOLD)
BODY = S("Normal", fontSize=9,  leading=13, fontName="Helvetica", textColor=C_DARK)
SM   = S("Normal", fontSize=7.5,leading=11, fontName="Helvetica", textColor=colors.HexColor("#6B7280"))
ITA  = S("Normal", fontSize=8.5,leading=12, fontName="Helvetica-Oblique",
          textColor=colors.HexColor("#6B7280"))

def hr(c=C_GOLD, t=1):
    return HRFlowable(width="100%", thickness=t, color=c, spaceAfter=5, spaceBefore=2)

def sec(title, sub=None):
    out = [Spacer(1, 0.2*cm), Paragraph(title, H2)]
    if sub: out.append(Paragraph(sub, ITA))
    out.append(hr())
    return out

def bul(text, c=C_TEAL):
    return Paragraph(f'<font color="#{c.hexval()[2:]}">▸</font>  {text}', BODY)

def tbl(data, col_widths, row_bgs=None, hdr_bg=C_DARK, tot_bg=C_GRN):
    t = Table(data, colWidths=col_widths)
    cmds = [
        ("BACKGROUND",(0,0),(-1,0), hdr_bg),
        ("TOPPADDING",(0,0),(-1,-1),4),("BOTTOMPADDING",(0,0),(-1,-1),4),
        ("LEFTPADDING",(0,0),(-1,-1),5),("GRID",(0,0),(-1,-1),0.25,C_DK2),
        ("VALIGN",(0,0),(-1,-1),"TOP"),
    ]
    if row_bgs:
        cmds.append(("ROWBACKGROUNDS",(0,1),(-1,-1),row_bgs))
    t.setStyle(TableStyle(cmds))
    return t

story = []

# COVER
story.append(Spacer(1,1.2*cm))
cov_data = [
    [Paragraph("<b>TITO AI — @TitoAIPH</b>",
               S("Normal",fontSize=30,fontName="Helvetica-Bold",textColor=C_GOLD,leading=36))],
    [Paragraph("Content Plan &amp; Boost Budget · W26–W28 · June–July 2026",
               S("Normal",fontSize=15,fontName="Helvetica-Bold",textColor=C_WHITE,leading=19))],
    [Paragraph("3 weeks · 9 posts · Php 6,000 total · Launch: Jun 22, 2026",
               S("Normal",fontSize=10,fontName="Helvetica",textColor=C_LG,leading=13))],
    [Spacer(1,0.3*cm)],
    [Paragraph("Progressive Ramp: Week 1 TikTok → Week 2 TikTok + Instagram → Week 3 All 3 Platforms",
               S("Normal",fontSize=9,fontName="Helvetica-Oblique",textColor=C_AMBE,leading=12))],
    [Spacer(1,0.3*cm)],
    [Paragraph('"The warm, relatable Tito who teaches everyday Filipinos to use AI for free — before it replaces their job."',
               S("Normal",fontSize=9,fontName="Helvetica-Oblique",textColor=C_AMBE,leading=12))],
    [Spacer(1,0.3*cm)],
    [Paragraph("Prepared by Jeff de las Armas  ·  June 2026",
               S("Normal",fontSize=8,fontName="Helvetica",textColor=C_GRAY,leading=11))],
]
ct = Table(cov_data, colWidths=[15*cm])
ct.setStyle(TableStyle([
    ("BACKGROUND",(0,0),(-1,-1),C_NAVY),
    ("TOPPADDING",(0,0),(-1,-1),6),("BOTTOMPADDING",(0,0),(-1,-1),6),
    ("LEFTPADDING",(0,0),(-1,-1),16),("RIGHTPADDING",(0,0),(-1,-1),16),
    ("LINEBELOW",(0,1),(-1,1),2,C_GOLD),
]))
story.append(ct)
story.append(Spacer(1,0.7*cm))

# TARGET MARKET
story += sec("Target Market — Mga Pamangkin", "Who we reach · why they follow · why they share")
seg_data = [
    [Paragraph(h, S("Normal",fontSize=8,fontName="Helvetica-Bold",textColor=C_WHITE,leading=11))
     for h in ["Audience Segment", "Size", "Fear / Need", "Share Trigger"]],
    *[
        [Paragraph(f'<font color="#{c.hexval()[2:]}">{seg}</font>',
                   S("Normal",fontSize=9,fontName="Helvetica-Bold",textColor=c,leading=12)),
         Paragraph(size, BODY),
         Paragraph(fn, BODY),
         Paragraph(st, BODY)]
        for seg, c, size, fn, st in [
            ("Freelancers",     C_GOLD,  "1.5M+ in PH",
             "Clients will use AI instead",
             "Posts with job titles + real Php outputs"),
            ("Guro / Teachers", C_TEAL,  "900K+ DepEd",
             "Curriculum disruption; admin overload",
             "Lesson plans in 2 mins — one guro shares to 40 families"),
            ("BPO Workers",     C_TIK,   "1.9M employees",
             "Automation layoffs; need to upskill NOW",
             "Fear + solution in same video"),
            ("Nanays / Tatays", C_PINK,  "Millions of parents",
             "Being left behind; no time to learn",
             "Family budget demos; cooking/homework use cases"),
            ("SME Owners",      C_FBB,   "14.9% use AI tools",
             "Losing to AI-enabled competition",
             "Real business wins with Php budgets"),
        ]
    ],
]
story.append(tbl(seg_data, [3.0*cm, 2.2*cm, 4.3*cm, 5.5*cm],
                 row_bgs=[C_DK2, C_DARK]*3))
story.append(Spacer(1,0.15*cm))

profile_data = [
    [Paragraph(h, S("Normal",fontSize=8,fontName="Helvetica-Bold",textColor=C_WHITE,leading=11))
     for h in ["Profile Attribute", "Detail"]],
    *[
        [Paragraph(k, S("Normal",fontSize=8.5,fontName="Helvetica-Bold",textColor=C_TEAL,leading=11)),
         Paragraph(v, BODY)]
        for k, v in [
            ("Age range",      "22–45 · Philippines · Taglish-speaking"),
            ("Pain point",     "Afraid AI will replace them before they learn to use it"),
            ("Motivation",     "Practical wins · libre tools · no jargon · real Filipino examples"),
            ("Watch trigger",  "Hook names their job or fear in the first 2 seconds"),
            ("Share trigger",  "Content that says 'this is for me' — specific personas and Php amounts"),
            ("Market gap",     "12.7M Filipinos exposed to AI disruption (ILO 2026) — almost none have a Taglish guide"),
        ]
    ],
]
story.append(tbl(profile_data, [3.5*cm, 11.5*cm],
                 row_bgs=[C_DK2, C_DARK]*4))
story.append(Spacer(1,0.4*cm))

# STRATEGY OVERVIEW
story += sec("Strategy Overview", "Php 2,000/week · 3 posts/week · launch Jun 22, 2026")
ramp_data_pdf = [
    [Paragraph(h, S("Normal",fontSize=8,fontName="Helvetica-Bold",textColor=C_WHITE,leading=11))
     for h in ["Week", "Dates", "Platforms", "Focus / Why", "Budget"]],
    *[
        [Paragraph(f'<font color="#{c.hexval()[2:]}">{wk}</font>',
                   S("Normal",fontSize=9,fontName="Helvetica-Bold",textColor=c,leading=12)),
         Paragraph(dates, BODY),
         Paragraph(f'<font color="#{c.hexval()[2:]}">{plats}</font>',
                   S("Normal",fontSize=9,fontName="Helvetica-Bold",textColor=c,leading=12)),
         Paragraph(why, BODY),
         Paragraph(bgt, S("Normal",fontSize=9,fontName="Helvetica-Bold",textColor=C_LGRN,leading=12))]
        for wk, dates, plats, why, bgt, c in [
            ("W26","Jun 22–28","TikTok only",
             "Learn TikTok Promote · establish baseline CPV + CPF · lowest risk","Php 2,000",C_TIK),
            ("W27","Jun 29–Jul 5","TikTok + Instagram",
             "Add IG Boost · apply W26 data · compare cost-per-follow across 2 platforms","Php 2,000",C_IGP),
            ("W28","Jul 6–12","TikTok + IG + Facebook",
             "Full 3-platform push · identify winner for July scale-up · FB Page required","Php 2,000",C_FBB),
        ]
    ],
    [Paragraph("TOTAL", S("Normal",fontSize=9,fontName="Helvetica-Bold",textColor=C_WHITE,leading=12)),
     Paragraph("3 weeks", BODY),
     Paragraph("3 platforms", BODY),
     Paragraph("Progressive · data-driven scale-up", BODY),
     Paragraph("Php 6,000", S("Normal",fontSize=10,fontName="Helvetica-Bold",textColor=C_AMBE,leading=13))],
]
story.append(tbl(ramp_data_pdf, [1.6*cm, 2.4*cm, 3.5*cm, 5.5*cm, 2.0*cm],
                 row_bgs=[C_DK2, C_DARK, C_DK2], tot_bg=C_GRN))
story.append(Spacer(1,0.4*cm))

# 9-POST CALENDAR
story += sec("3-Week Content Calendar", "9 posts · Mon/Wed/Fri · drop 7–8 PM PHT")
cal_data_pdf = [
    [Paragraph(h, S("Normal",fontSize=7.5,fontName="Helvetica-Bold",textColor=C_WHITE,leading=10))
     for h in ["Wk", "Date", "Day", "Topic", "Platform", "Budget", "Boost Dates"]],
    *[
        [Paragraph(f'<font color="#{wc.hexval()[2:]}">{wk}</font>',
                   S("Normal",fontSize=8.5,fontName="Helvetica-Bold",leading=11)),
         Paragraph(date, BODY),
         Paragraph(day, S("Normal",fontSize=8.5,fontName="Helvetica-Bold",textColor=C_GOLD,leading=11)),
         Paragraph(topic, BODY),
         Paragraph(f'<font color="#{pc.hexval()[2:]}">{plat}</font>',
                   S("Normal",fontSize=8.5,fontName="Helvetica-Bold",leading=11)),
         Paragraph(bgt, S("Normal",fontSize=9,fontName="Helvetica-Bold",textColor=C_LGRN,leading=12)),
         Paragraph(bd, SM)]
        for wk,wc,date,day,topic,plat,pc,bgt,bd in [
            ("W26",C_TIK,"Jun 22","MON","Ang Sabi Nila: Pang-Matalino Lang Iyan","TikTok",C_TIK,"Php 700","Jun 23–24"),
            ("W26",C_TIK,"Jun 25","WED","Gemini para sa Guro","TikTok",C_TIK,"Php 700","Jun 26–27"),
            ("W26",C_TIK,"Jun 27","FRI","1 Buwan Kasama Tito AI — Ang Resulta","TikTok",C_TIK,"Php 600","Jun 28–29"),
            ("W27",C_IGP,"Jun 30","MON","July Teaser — Ano ang Susunod?","TikTok",C_TIK,"Php 700","Jul 1–2"),
            ("W27",C_IGP,"Jul 2","WED","TBD — Wednesday Demo","Instagram",C_IGP,"Php 700","Jul 3–4"),
            ("W27",C_IGP,"Jul 4","FRI","TBD — Friday Story","TikTok",C_TIK,"Php 600","Jul 5–6"),
            ("W28",C_FBB,"Jul 7","MON","TBD — Monday AI Tip","TikTok",C_TIK,"Php 600","Jul 8–9"),
            ("W28",C_FBB,"Jul 9","WED","TBD — Wednesday Demo","Facebook",C_FBB,"Php 800","Jul 10–11"),
            ("W28",C_FBB,"Jul 11","FRI","TBD — Friday Story","Instagram",C_IGP,"Php 600","Jul 12–13"),
        ]
    ],
    [Paragraph("TOTAL",S("Normal",fontSize=9,fontName="Helvetica-Bold",textColor=C_WHITE,leading=12)),
     Paragraph("—",BODY), Paragraph("9 posts",BODY),
     Paragraph("3 weeks · 3 platforms",BODY),
     Paragraph("—",BODY),
     Paragraph("Php 6,000",S("Normal",fontSize=10,fontName="Helvetica-Bold",textColor=C_AMBE,leading=13)),
     Paragraph("—",BODY)],
]
story.append(tbl(cal_data_pdf, [1.0*cm,1.6*cm,1.0*cm,4.5*cm,2.2*cm,1.5*cm,2.2*cm],
                 row_bgs=[C_DK2,C_DARK]*5))
story.append(Spacer(1,0.4*cm))

# BUDGET BREAKDOWN BY WEEK
story += sec("Budget Breakdown by Week")
for wk_label, wk_col, wk_rows, wk_total, note in [
    ("WEEK 1 — W26 (Jun 22–28) · TikTok Only", C_TIK, [
        ("Mon Jun 22","Ang Sabi Nila: Pang-Matalino Lang Iyan","TikTok Promote","Php 350/day × 2","Php 700","Jun 23–24"),
        ("Wed Jun 25","Gemini para sa Guro",                    "TikTok Promote","Php 350/day × 2","Php 700","Jun 26–27"),
        ("Fri Jun 27","1 Buwan Kasama Tito AI — Ang Resulta",   "TikTok Promote","Php 300/day × 2","Php 600","Jun 28–29"),
    ], "Php 2,000", "Goal: Video views · Automatic audience · Start morning after posting"),
    ("WEEK 2 — W27 (Jun 29–Jul 5) · TikTok + Instagram", C_IGP, [
        ("Mon Jun 30","July Teaser — Ano ang Susunod?",  "TikTok Promote","Php 350/day × 2","Php 700","Jul 1–2"),
        ("Wed Jul 2", "TBD — Wednesday Demo",            "Instagram Boost","Php 350/day × 2","Php 700","Jul 3–4"),
        ("Fri Jul 4", "TBD — Friday Story",              "TikTok Promote","Php 300/day × 2","Php 600","Jul 5–6"),
    ], "Php 2,000", "Apply W26 learnings · TikTok: Video views · Instagram: Profile visits"),
    ("WEEK 3 — W28 (Jul 6–12) · All 3 Platforms", C_FBB, [
        ("Mon Jul 7","TBD — Monday AI Tip",  "TikTok Promote","Php 300/day × 2","Php 600","Jul 8–9"),
        ("Wed Jul 9","TBD — Wednesday Demo", "Facebook Boost","Php 400/day × 2","Php 800","Jul 10–11"),
        ("Fri Jul 11","TBD — Friday Story",  "Instagram Boost","Php 300/day × 2","Php 600","Jul 12–13"),
    ], "Php 2,000", "⚠️ Facebook Page required before boosting — convert profiles at facebook.com/pages/create"),
]:
    story.append(KeepTogether([
        Paragraph(wk_label, S("Normal",fontSize=10,fontName="Helvetica-Bold",
                               textColor=wk_col,leading=13)),
        Paragraph(note, ITA),
    ]))
    wk_data_pdf = [
        [Paragraph(h, S("Normal",fontSize=7.5,fontName="Helvetica-Bold",textColor=C_WHITE,leading=10))
         for h in ["Day","Post","Platform","Rate","Total","Boost Dates"]],
        *[
            [Paragraph(d,BODY),Paragraph(t,BODY),
             Paragraph(f'<font color="#{wk_col.hexval()[2:]}">{p}</font>',
                       S("Normal",fontSize=8.5,fontName="Helvetica-Bold",leading=11)),
             Paragraph(r,SM),
             Paragraph(tot,S("Normal",fontSize=9,fontName="Helvetica-Bold",textColor=C_LGRN,leading=12)),
             Paragraph(bd,SM)]
            for d,t,p,r,tot,bd in wk_rows
        ],
        [Paragraph("",BODY),Paragraph("",BODY),Paragraph("",BODY),
         Paragraph("Week Total",S("Normal",fontSize=9,fontName="Helvetica-Bold",textColor=C_WHITE,leading=12)),
         Paragraph(wk_total,S("Normal",fontSize=10,fontName="Helvetica-Bold",textColor=C_AMBE,leading=13)),
         Paragraph("",BODY)],
    ]
    wt = Table(wk_data_pdf, colWidths=[2.0*cm,4.5*cm,2.8*cm,2.4*cm,1.5*cm,2.3*cm])
    wt.setStyle(TableStyle([
        ("BACKGROUND",(0,0),(-1,0),C_DARK),
        ("BACKGROUND",(0,len(wk_rows)+1),(-1,len(wk_rows)+1),C_GRN),
        ("ROWBACKGROUNDS",(0,1),(-1,len(wk_rows)),[C_DK2,C_DARK]*3),
        ("TOPPADDING",(0,0),(-1,-1),4),("BOTTOMPADDING",(0,0),(-1,-1),4),
        ("LEFTPADDING",(0,0),(-1,-1),5),("GRID",(0,0),(-1,-1),0.25,C_DK2),
        ("VALIGN",(0,0),(-1,-1),"TOP"),
    ]))
    story.append(wt)
    story.append(Spacer(1,0.3*cm))

# KPIs
story += sec("KPIs to Track Weekly", "Pull every Sunday · compare across weeks · share with Jeff")
kpi_pdf_data = [
    [Paragraph(h, S("Normal",fontSize=8,fontName="Helvetica-Bold",textColor=C_WHITE,leading=11))
     for h in ["Metric", "Platform", "W1 Target", "W3 Target"]],
    *[
        [Paragraph(m,BODY),
         Paragraph(f'<font color="#{pc.hexval()[2:]}">{p}</font>',
                   S("Normal",fontSize=8.5,fontName="Helvetica-Bold",leading=11)),
         Paragraph(t1,S("Normal",fontSize=9,fontName="Helvetica-Bold",textColor=C_LGRN,leading=12)),
         Paragraph(t3,S("Normal",fontSize=9,fontName="Helvetica-Bold",textColor=C_GOLD,leading=12))]
        for m,p,pc,t1,t3 in [
            ("Reach per boosted post","TikTok",C_TIK,"3,000+","7,000+"),
            ("Video views per post","TikTok",C_TIK,"5,000+","15,000+"),
            ("New followers (3 posts)","TikTok",C_TIK,"60–120","250+"),
            ("Cost per follow","TikTok",C_TIK,"< Php 25","< Php 15"),
            ("Reach per boost","Instagram",C_IGP,"2,000+","5,000+"),
            ("New followers (per post)","Instagram",C_IGP,"20–40","100+"),
            ("Cost per follow","Instagram",C_IGP,"< Php 25","< Php 15"),
            ("Reach per boost","Facebook",C_FBB,"2,500+","6,000+"),
            ("New Page likes","Facebook",C_FBB,"20–40","100+"),
            ("Cost per Page like","Facebook",C_FBB,"< Php 25","< Php 15"),
        ]
    ],
]
story.append(tbl(kpi_pdf_data, [5.0*cm,2.8*cm,2.8*cm,4.4*cm],
                 row_bgs=[C_DK2,C_DARK]*6))
story.append(Spacer(1,0.3*cm))

# Scale rules
story.append(Paragraph("Scale Rules — Week 4+ (July)", H3))
for rule in [
    "50+ followers from a single boost → extend that post 2 more days at the same daily rate",
    "Cost-per-follow drops below Php 12 → increase that platform's weekly budget by Php 400",
    "Platform with lowest CPF across W26–W28 → gets the biggest share of July budget",
    "Post earns < 20 followers after boost → swap to a different platform or content slot next week",
]:
    story.append(bul(rule, C_GOLD))

story.append(Spacer(1,0.8*cm))
story.append(Paragraph(
    "Tito AI @TitoAIPH  ·  Content Plan &amp; Boost Budget W26–W28  ·  Prepared by Jeff de las Armas  ·  June 2026",
    S("Normal",fontSize=7,fontName="Helvetica",textColor=C_GRAY,leading=10,alignment=TA_CENTER)))

doc = SimpleDocTemplate(OUT_PDF, pagesize=A4,
                        leftMargin=1*cm, rightMargin=1*cm,
                        topMargin=1.2*cm, bottomMargin=1.2*cm,
                        title="Tito AI Content Plan Boost Budget W26-W28",
                        author="Jeff de las Armas")
doc.build(story)
print(f"✓ PDF   → {OUT_PDF}")
